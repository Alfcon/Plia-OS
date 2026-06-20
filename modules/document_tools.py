from __future__ import annotations
from core.registry import tool


def _expand(path: str) -> str:
    import os
    return os.path.expanduser(path)


@tool("Read text from a PDF file. max_pages=0 means all pages.")
def read_pdf(path: str, max_pages: int = 0) -> str:
    try:
        import pypdf
    except ImportError:
        return "pypdf not installed. Run: pip install pypdf"
    p = _expand(path)
    try:
        reader = pypdf.PdfReader(p)
    except FileNotFoundError:
        return f"File not found: {path}"
    except Exception as exc:
        return f"Error opening PDF: {exc}"
    pages = reader.pages if not max_pages else reader.pages[:max_pages]
    parts = []
    for i, page in enumerate(pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"[Page {i}]\n{text.strip()}")
    if not parts:
        return "No extractable text found in PDF."
    return "\n\n".join(parts)


@tool("Read text from a Word document (.docx).")
def read_docx(path: str) -> str:
    try:
        import docx
    except ImportError:
        return "python-docx not installed. Run: pip install python-docx"
    p = _expand(path)
    try:
        doc = docx.Document(p)
    except FileNotFoundError:
        return f"File not found: {path}"
    except Exception as exc:
        return f"Error opening DOCX: {exc}"
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    if not paragraphs:
        return "No text found in document."
    return "\n".join(paragraphs)


@tool("Read an Excel spreadsheet (.xlsx). Returns all sheets as CSV-style text, or a named sheet.")
def read_xlsx(path: str, sheet: str = "") -> str:
    try:
        import openpyxl
    except ImportError:
        return "openpyxl not installed. Run: pip install openpyxl"
    p = _expand(path)
    try:
        wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
    except FileNotFoundError:
        return f"File not found: {path}"
    except Exception as exc:
        return f"Error opening XLSX: {exc}"

    sheet_names = sheet.split(",") if sheet else wb.sheetnames
    parts = []
    for name in sheet_names:
        name = name.strip()
        if name not in wb.sheetnames:
            parts.append(f"[{name}] Sheet not found. Available: {', '.join(wb.sheetnames)}")
            continue
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        parts.append(f"[{name}]\n" + ("\n".join(rows) if rows else "(empty)"))
    wb.close()
    return "\n\n".join(parts)


@tool("Read text from a PowerPoint presentation (.pptx). Returns slide titles and content.")
def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "python-pptx not installed. Run: pip install python-pptx"
    p = _expand(path)
    try:
        prs = Presentation(p)
    except FileNotFoundError:
        return f"File not found: {path}"
    except Exception as exc:
        return f"Error opening PPTX: {exc}"
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    if not parts:
        return "No text found in presentation."
    return "\n\n".join(parts)
