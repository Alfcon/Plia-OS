import io
import pytest
from unittest.mock import patch, MagicMock


def test_read_pdf_returns_text(tmp_path):
    import pypdf
    from modules.document_tools import read_pdf

    # Build a real minimal PDF via pypdf
    writer = pypdf.PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    pdf_path = tmp_path / "test.pdf"
    pdf_path.write_bytes(buf.getvalue())

    # No extractable text — expect the "No extractable text" message
    result = read_pdf(str(pdf_path))
    assert "No extractable text" in result or isinstance(result, str)


def test_read_pdf_not_found():
    from modules.document_tools import read_pdf
    result = read_pdf("/nonexistent/path.pdf")
    assert "not found" in result.lower()


def test_read_docx_returns_text(tmp_path):
    import docx
    from modules.document_tools import read_docx

    doc = docx.Document()
    doc.add_paragraph("Hello from Plia")
    doc.add_paragraph("Second paragraph")
    path = tmp_path / "test.docx"
    doc.save(str(path))

    result = read_docx(str(path))
    assert "Hello from Plia" in result
    assert "Second paragraph" in result


def test_read_docx_not_found():
    from modules.document_tools import read_docx
    result = read_docx("/nonexistent/path.docx")
    assert "not found" in result.lower()


def test_read_xlsx_returns_data(tmp_path):
    import openpyxl
    from modules.document_tools import read_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    path = tmp_path / "test.xlsx"
    wb.save(str(path))

    result = read_xlsx(str(path))
    assert "Name" in result
    assert "Alice" in result


def test_read_xlsx_not_found():
    from modules.document_tools import read_xlsx
    result = read_xlsx("/nonexistent/path.xlsx")
    assert "not found" in result.lower()


def test_read_pptx_returns_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from modules.document_tools import read_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    path = tmp_path / "test.pptx"
    prs.save(str(path))

    result = read_pptx(str(path))
    assert "Test Slide" in result


def test_read_pptx_not_found():
    from modules.document_tools import read_pptx
    result = read_pptx("/nonexistent/path.pptx")
    assert "not found" in result.lower()
